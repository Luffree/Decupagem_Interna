"""
PPTX content extractor using python-pptx.
Extracts text, tables, shapes, notes, hyperlinks and visual elements per slide.
"""
from dataclasses import dataclass, field
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class SlideContent:
    slide_number: int
    title: str
    text_blocks: list[str]
    tables: list[list[list[str]]]
    shapes: list[dict]
    notes: str
    hyperlinks: list[str]
    placeholder_data: list[dict]


@dataclass
class PPTXContent:
    filename: str
    total_slides: int
    slides: list[SlideContent] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)


def _extract_text_from_shape(shape) -> str:
    """Recursively extract all text from a shape including grouped shapes."""
    texts = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            t = _extract_text_from_shape(s)
            if t:
                texts.append(t)
    elif shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            para_text = "".join(run.text for run in para.runs).strip()
            if para_text:
                texts.append(para_text)
    return "\n".join(texts)


def _extract_hyperlinks(shape) -> list[str]:
    links = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            links.extend(_extract_hyperlinks(s))
    elif shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.hyperlink and run.hyperlink.address:
                    links.append(run.hyperlink.address)
    if hasattr(shape, "click_action") and shape.click_action:
        try:
            if shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                links.append(shape.click_action.hyperlink.address)
        except Exception:
            pass
    return links


def _extract_table(shape) -> list[list[str]]:
    rows = []
    for row in shape.table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip())
        rows.append(cells)
    return rows


def _get_shape_info(shape) -> dict:
    info = {
        "name": shape.name,
        "type": str(shape.shape_type),
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
        "text": "",
        "fill_color": None,
        "line_color": None,
    }
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return info

    try:
        if shape.has_text_frame:
            info["text"] = shape.text_frame.text.strip()
    except Exception:
        pass

    try:
        fill = shape.fill
        if fill.type is not None:
            if hasattr(fill, "fore_color") and fill.fore_color and fill.fore_color.type is not None:
                try:
                    info["fill_color"] = str(fill.fore_color.rgb)
                except Exception:
                    pass
    except Exception:
        pass

    try:
        line = shape.line
        if line and line.color and line.color.type is not None:
            try:
                info["line_color"] = str(line.color.rgb)
            except Exception:
                pass
    except Exception:
        pass

    return info


def extract_pptx(file_path: str) -> PPTXContent:
    prs = Presentation(file_path)

    content = PPTXContent(
        filename=file_path,
        total_slides=len(prs.slides),
        slides=[],
        metadata={
            "slide_width_emu": prs.slide_width,
            "slide_height_emu": prs.slide_height,
        },
    )

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        title = ""
        text_blocks = []
        tables = []
        shapes_info = []
        hyperlinks = []

        # Extract title from layout or placeholder
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            # Collect shape metadata
            shapes_info.append(_get_shape_info(shape))

            # Tables
            if shape.has_table:
                tables.append(_extract_table(shape))
                continue

            # Hyperlinks
            hyperlinks.extend(_extract_hyperlinks(shape))

            # Text
            text = _extract_text_from_shape(shape)
            if text and text != title:
                text_blocks.append(text)

        # Placeholder breakdown for more context
        placeholder_data = []
        for ph in slide.placeholders:
            try:
                ph_info = {
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type),
                    "text": ph.text.strip() if ph.has_text_frame else "",
                }
                placeholder_data.append(ph_info)
            except Exception:
                pass

        # Notes
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                notes_text = notes_tf.text.strip()

        content.slides.append(SlideContent(
            slide_number=slide_num,
            title=title,
            text_blocks=text_blocks,
            tables=tables,
            shapes=shapes_info,
            notes=notes_text,
            hyperlinks=list(set(hyperlinks)),
            placeholder_data=placeholder_data,
        ))

    return content


def extract_pptx_bytes(file_bytes: bytes, filename: str) -> PPTXContent:
    """Extract content from PPTX given as bytes."""
    import tempfile, os
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        result = extract_pptx(tmp_path)
        result.filename = filename
        return result
    finally:
        os.unlink(tmp_path)


def pptx_content_to_dict(content: PPTXContent) -> dict:
    """Convert PPTXContent to a JSON-serializable dict for the analyzer."""
    return {
        "filename": content.filename,
        "total_slides": content.total_slides,
        "metadata": {k: str(v) for k, v in content.metadata.items()},
        "slides": [
            {
                "slide_number": s.slide_number,
                "title": s.title,
                "text_blocks": s.text_blocks,
                "tables": s.tables,
                "notes": s.notes,
                "hyperlinks": s.hyperlinks,
                "placeholder_data": s.placeholder_data,
                "shapes_summary": [
                    {
                        "name": sh["name"],
                        "type": sh["type"],
                        "text": sh["text"],
                        "fill_color": sh["fill_color"],
                        "line_color": sh["line_color"],
                    }
                    for sh in s.shapes
                    if sh.get("text") or sh.get("fill_color")
                ],
            }
            for s in content.slides
        ],
    }
